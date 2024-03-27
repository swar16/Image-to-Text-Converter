import os
import smtplib
import time
from email.mime.application import MIMEApplication
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
import cv2
import numpy as np


def binarize_image(image):
    _, binary = cv2.threshold(image, 128, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)
    return binary


def apply_morphology(image):
    kernel = np.ones((2, 2), np.uint8)
    morph = cv2.morphologyEx(image, cv2.MORPH_CLOSE, kernel)
    return morph


def remove_shadows(image):
    if len(image.shape) > 2 and image.shape[2] == 3:

        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    else:

        gray = image

    dilated_img = cv2.dilate(gray, np.ones((7, 7), np.uint8))
    blurred_img = cv2.medianBlur(dilated_img, 21)
    divided = np.divide(gray, blurred_img)
    normalized_img = np.clip(divided * 255, 0, 255).astype(np.uint8)
    return normalized_img


def resize_image(image, height=900):
    aspect_ratio = image.shape[1] / image.shape[0]
    width = int(aspect_ratio * height)
    resized_img = cv2.resize(image, (width, height))
    return resized_img


def apply_adaptive_threshold(image):
    adaptive_thresh = cv2.adaptiveThreshold(image, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)
    return adaptive_thresh


def preprocess_image(image_path):
    img = cv2.imread(image_path)

    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

    denoised = cv2.fastNlMeansDenoising(gray, h=10, templateWindowSize=7, searchWindowSize=21)

    alpha = 1.5
    beta = 0
    adjusted = cv2.convertScaleAbs(denoised, alpha=alpha, beta=beta)

    binary_image = binarize_image(adjusted)

    morph_image = apply_morphology(binary_image)

    no_shadow_image = remove_shadows(morph_image)

    resized_image = resize_image(no_shadow_image)

    adaptive_thresh_image = apply_adaptive_threshold(resized_image)

    kernel = np.array([[-1, -1, -1],
                       [-1, 9, -1],
                       [-1, -1, -1]])
    sharpened = cv2.filter2D(adaptive_thresh_image, -1, kernel)

    coords = np.column_stack(np.where(sharpened > 0))
    angle = cv2.minAreaRect(coords)[-1]
    if angle < -45:
        angle = -(90 + angle)
    else:
        angle = -angle
    (h, w) = sharpened.shape[:2]
    center = (w // 2, h // 2)
    M = cv2.getRotationMatrix2D(center, angle, 1.0)
    deskewed = cv2.warpAffine(sharpened, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)

    return deskewed


def getImgsFromPDF(pdf_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)

    pdf_document = fitz.open(pdf_path)

    for page_index in range(len(pdf_document)):
        page = pdf_document[page_index]

        image_list = page.get_images(full=True)

        if image_list:
            print(f"[+] Found a total of {len(image_list)} images in page {page_index}")
        else:
            print("[!] No images found on page", page_index)

        for image_index, img in enumerate(image_list):
            xref = img[0]

            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]

            image_ext = base_image["ext"]

            image_name = f"image{page_index + 1}_{image_index + 1}.{image_ext}"

            image_path = os.path.join(output_folder, image_name)
            with open(image_path, "wb") as image_file:
                image_file.write(image_bytes)

            print(f"Saved image {image_index + 1} from page {page_index + 1}")

    pdf_document.close()


def getTextFromImgs(images_folder, text_folder):
    # Ensure output directory exists
    os.makedirs(text_folder, exist_ok=True)

    for image_name in os.listdir(images_folder):
        if image_name.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif")):
            image_path = os.path.join(images_folder, image_name)

            processed_image = preprocess_image(image_path)

            image = Image.fromarray(processed_image)

            text = pytesseract.image_to_string(image, lang='eng')

            text_file_name = os.path.splitext(image_name)[0] + '.txt'
            text_file_path = os.path.join(text_folder, text_file_name)

            with open(text_file_path, 'w', encoding='utf-8') as text_file:
                text_file.write(text)

            print(f"Processed {image_name}")


def getPPTFromImgText(images_folder, ppt_path):
    prs = Presentation()
    for image_name in sorted(os.listdir(images_folder)):
        if image_name.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif")):
            image_path = os.path.join(images_folder, image_name)

            image = Image.open(image_path)
            text = pytesseract.image_to_string(image, lang='eng')

            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)

            txBox = slide.shapes.add_textbox(Pt(50), Pt(50), prs.slide_width - Pt(100), prs.slide_height - Pt(100))
            tf = txBox.text_frame
            tf.text = text

    prs.save(ppt_path)
    print(f"PowerPoint saved at {ppt_path}")


def send_email_with_ppt(recipient_email, pptx_file_path):
    time.sleep(3)

    HOST = "smtp-mail.outlook.com"
    PORT = 587
    FROM_EMAIL = "codesteinsprojectmail@gmail.com"
    TO_EMAIL = recipient_email
    PASSWORD = "codes#2oo345"

    msg = MIMEMultipart()
    msg['From'] = FROM_EMAIL
    msg['To'] = TO_EMAIL
    msg['Subject'] = "Here is the required Powerpoint Presentation"

    message = "Please find the attached Powerpoint presentation."
    msg.attach(MIMEText(message, 'plain'))

    with open(pptx_file_path, 'rb') as file:
        pptx_attachment = MIMEApplication(file.read(), _subtype="pptx")
    pptx_attachment.add_header('Content-Disposition', f'attachment; filename={pptx_file_path}')
    msg.attach(pptx_attachment)

    try:
        smtp = smtplib.SMTP(HOST, PORT)
        smtp.starttls()
        smtp.login(FROM_EMAIL, PASSWORD)
        smtp.sendmail(FROM_EMAIL, TO_EMAIL, msg.as_string())
        smtp.quit()
        print("Email sent successfully!")

    except Exception as e:
        print(f"Email could not be sent. Error: {e}")


def send_feedback_email(recipient_email, feedback_text):
    HOST = "smtp-mail.outlook.com"
    PORT = 587
    FROM_EMAIL = "codesteinsprojectmail@gmail.com"
    PASSWORD = "codes#2oo345"

    msg = MIMEMultipart('alternative')
    msg['From'] = FROM_EMAIL
    msg['To'] = recipient_email
    msg['Subject'] = "New Feedback Submission"

    body_content = """
    <html>
    <head>
        <style>
            body {{
                font-family: Arial, sans-serif;
                margin: 20px;
                background-color: #f4f4f4;
                color: #555;
            }}
            .container {{
                background-color: #fff;
                padding: 20px;
                border-radius: 5px;
                box-shadow: 0 0 10px rgba(0, 0, 0, 0.1);
            }}
            .header {{
                font-size: 24px;
                font-weight: bold;
                margin-bottom: 20px;
            }}
            .feedback-text {{
                margin-top: 20px;
                line-height: 1.6;
            }}
            .footer {{
                margin-top: 30px;
                font-size: 12px;
                text-align: center;
            }}
        </style>
    </head>
    <body>
        <div class="container">
            <div class="header">Feedback Received</div>
            <div class="feedback-text">
                <p><b>Feedback:</b></p>
                <p>{feedback_text}</p>
            </div>
            <div class="footer">
                This is an automated message. Please do not reply directly to this email.
            </div>
        </div>
    </body>
    </html>
    """.format(feedback_text=feedback_text)

    part2 = MIMEText(body_content, 'html')
    msg.attach(part2)

    try:
        smtp = smtplib.SMTP(HOST, PORT)
        smtp.starttls()
        smtp.login(FROM_EMAIL, PASSWORD)
        smtp.sendmail(FROM_EMAIL, recipient_email, msg.as_string())
        smtp.quit()
        print("Feedback email sent successfully!")

    except Exception as e:
        print(f"Feedback email could not be sent. Error: {e}")
